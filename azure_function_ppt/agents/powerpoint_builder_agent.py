"""
PowerPoint Builder Agent - With template support
"""
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from typing import Dict, Any, Optional
from agents.core.base_agent import BaseAgent
from config import get_template_path
import json
import io
import os
import re

class PowerPointBuilderAgent(BaseAgent):
    """Builds PowerPoint files with 16:9 aspect ratio and simple 2-color theme"""

    # Simple color theme
    PRIMARY_COLOR = RGBColor(88, 77, 193)    # #584dc1 (purple)
    ACCENT_COLOR = RGBColor(209, 185, 91)    # #d1b95b (gold)
    TEXT_COLOR = RGBColor(51, 51, 51)        # Dark gray

    agent_description = "PowerPoint file generation with 16:9 aspect ratio and template support"
    agent_use_cases = [
        "PowerPoint file creation from slide content",
        "16:9 widescreen presentation generation", 
        "Template-based presentation generation",
        "Custom branding and layouts", 
        "Professional formatting"
    ]

    def __init__(self, **kwargs):
        super().__init__()
    
    def _set_16_9_aspect_ratio(self, prs: Presentation):
        """Set presentation to 16:9 aspect ratio - STRICTLY ENFORCED"""
        prs.slide_width = Inches(16)
        prs.slide_height = Inches(9)
        print(f"Enforced 16:9 aspect ratio: {prs.slide_width.inches:.1f}\" x {prs.slide_height.inches:.1f}\"")

    async def process(self, slide_content: str, context_metadata: Optional[Dict[str, Any]] = None) -> bytes:
        """Generate PowerPoint file from slide content with template support"""
        try:
            slides_data = self._parse_slide_content(slide_content)
            
            # Use template if available (controlled by config.py)
            template_path = get_template_path("default")
            if template_path:
                print(f"Using template: {template_path}")
                prs = Presentation(template_path)
                print(f"Template original dimensions: {prs.slide_width.inches:.1f}\" x {prs.slide_height.inches:.1f}\"")
                # ENFORCE 16:9 even with template for consistency
                self._set_16_9_aspect_ratio(prs)
            else:
                print("Using python-pptx default template (no custom design)")
                prs = Presentation()
                # Only set 16:9 when not using custom template
                self._set_16_9_aspect_ratio(prs)
            
            for slide_info in slides_data:
                self._create_slide(prs, slide_info)
            
            ppt_buffer = io.BytesIO()
            prs.save(ppt_buffer)
            ppt_buffer.seek(0)
            
            return ppt_buffer.read()
            
        except Exception as e:
            print(f"PowerPoint building error: {str(e)}")
            # If template failed, try without template
            if "template" in str(e).lower():
                print("Template failed, falling back to default")
                try:
                    slides_data = self._parse_slide_content(slide_content)
                    prs = Presentation()  # Use python-pptx default
                    
                    # Set to 16:9 aspect ratio for fallback
                    self._set_16_9_aspect_ratio(prs)
                    
                    for slide_info in slides_data:
                        self._create_slide(prs, slide_info)
                    
                    ppt_buffer = io.BytesIO()
                    prs.save(ppt_buffer)
                    ppt_buffer.seek(0)
                    
                    return ppt_buffer.read()
                except Exception as fallback_error:
                    raise Exception(f"Failed to generate PowerPoint even with fallback: {str(fallback_error)}")
            
            raise Exception(f"Failed to generate PowerPoint: {str(e)}")

    def _parse_slide_content(self, slide_content: str) -> list:
        """Parse slide content into structured data"""
        try:
            # Clean up potential markdown formatting around JSON
            if '```json' in slide_content:
                slide_content = slide_content.split('```json')[1].split('```')[0]

            if slide_content.strip().startswith(('[', '{')):
                content_data = json.loads(slide_content)
                # Handle cases where the JSON is a dict with a 'slides' key
                if isinstance(content_data, dict):
                    return content_data.get('slides') or content_data.get('presentation_structure', [])
                return content_data
        except (json.JSONDecodeError, IndexError):
            # Fallback to markdown if JSON parsing fails
            pass
        
        return self._parse_markdown_content(slide_content)

    def _parse_markdown_content(self, content: str) -> list:
        """Parse markdown-style slide content"""
        slides = []
        current_slide = None
        
        for line in content.split('\n'):
            line = line.strip()
            if not line:
                continue
                
            if line.startswith(('# ', '## ')):
                if current_slide:
                    slides.append(current_slide)
                current_slide = {
                    "title": line.lstrip('# '),
                    "content": [],
                    "layout": "TITLE_SLIDE" if not slides else "CONTENT_SLIDE"
                }
            elif line.startswith(('- ', '* ')) and current_slide:
                current_slide["content"].append(line.lstrip('- *'))
            elif line and current_slide:
                current_slide["content"].append(line)
        
        if current_slide:
            slides.append(current_slide)
        
        if not slides:
            slides = [{
                "title": "Generated Presentation",
                "content": ["Content extracted from document"],
                "layout": "TITLE_SLIDE"
            }]
        
        return slides

    def _create_slide(self, prs: Presentation, slide_info: dict):
        """Create individual slide with formatting based on slide type"""
        # Get slide type and determine layout
        slide_type = slide_info.get("slide_type") or slide_info.get("layout", "CONTENT_SLIDE")
        title = slide_info.get("title", "Slide Title")
        
        # Content can come from either "content" (detailed) or "content_outline" (structure)
        content = slide_info.get("content") or slide_info.get("content_outline", [])
        
        try:
            # Determine the appropriate PowerPoint layout based on slide type
            slide_layout_index = self._get_layout_index(slide_type, prs)
            slide_layout = prs.slide_layouts[slide_layout_index]
            slide = prs.slides.add_slide(slide_layout)
            
            print(f"Creating slide: {title} (type: {slide_type}, layout: {slide_layout_index})")
            print(f"Available placeholders: {len(slide.placeholders)}")
            
            # Apply type-specific formatting
            self._format_slide_by_type(slide, slide_type, title, content)
            
        except Exception as e:
            print(f"Error creating slide '{title}' of type '{slide_type}': {str(e)}")
            # Try with a simpler approach - just set title if possible
            slide_layout = prs.slide_layouts[0]  # Use title layout as fallback
            slide = prs.slides.add_slide(slide_layout)
            if slide.shapes.title:
                slide.shapes.title.text = title
            raise e

    def _get_layout_index(self, slide_type: str, prs: Presentation) -> int:
        """Get appropriate layout index - SIMPLE 3-LAYOUT SYSTEM"""
        max_layouts = len(prs.slide_layouts)
        print(f"Template has {max_layouts} layouts available")
        
        # SIMPLE 3-LAYOUT SYSTEM:
        # Layout 0: Opening slide (title/intro) - used once
        # Layout 1: Content slide (standard) - used for most slides
        # Layout 2: Ending slide (thank you/NCS) - used once
        
        if slide_type == "TITLE_SLIDE":
            layout_index = 0  # Opening slide
        elif slide_type == "THANK_YOU_SLIDE":
            layout_index = 2  # Ending slide
        else:
            layout_index = 1  # Content slide for everything else
        
        # Fallback if layout doesn't exist
        if layout_index >= max_layouts:
            print(f"Warning: Layout {layout_index} not available, using layout 0")
            layout_index = 0
            
        print(f"Using layout {layout_index} for slide type {slide_type}")
        return layout_index

    def _format_slide_by_type(self, slide, slide_type: str, title: str, content):
        """Apply type-specific formatting to slide"""
        try:
            # Set title - DON'T apply custom formatting, use template's formatting
            if slide.shapes.title:
                print(f"Found title shape, setting to: {title}")
                slide.shapes.title.text = title
                print(f"Successfully set slide title: {title}")
                # Remove custom formatting to preserve template fonts
                # self._apply_title_format_by_type(slide.shapes.title, slide_type)  # REMOVED
            else:
                print(f"No title shape found for slide with title: {title}")
        except Exception as e:
            print(f"Error setting title '{title}': {e}")
        
        try:
            # SIMPLIFIED: Use the same content formatting for all slide types
            # Only the layout (0, 1, or 2) determines the visual design
            self._format_any_slide_content(slide, content, slide_type)
        except Exception as e:
            print(f"Error formatting slide content for type {slide_type}: {e}")
            # Just leave the slide as-is if formatting fails

    def _format_title_slide(self, slide, content):
        """Format title slide with subtitle"""
        subtitle_placeholder = self._find_content_placeholder(slide)
        if subtitle_placeholder:
            try:
                subtitle_text = content[0] if content else "Professional Business Presentation"
                subtitle_placeholder.text = str(subtitle_text)
                self._apply_smart_text_formatting(subtitle_placeholder, str(subtitle_text))
                print(f"Set title slide subtitle: {subtitle_text}")
            except Exception as e:
                print(f"Error setting title slide content: {e}")
        else:
            print("No subtitle placeholder found for title slide")

    def _format_ncs_ending_slide(self, slide, content):
        """Format NCS Singapore ending slide - preserve template design"""
        content_placeholder = self._find_content_placeholder(slide)
        if content_placeholder:
            try:
                # FORCE: Always only "Thank you" regardless of AI-generated content
                ending_text = "Thank you"
                content_placeholder.text = ending_text
                self._apply_smart_text_formatting(content_placeholder, ending_text)
                print(f"Set NCS ending slide content: '{ending_text}'")
                
            except Exception as e:
                print(f"Error setting NCS ending content: {e}")
        else:
            print("No content placeholder found for NCS ending slide")

    def _format_standout_slide(self, slide, content):
        """Format stand out message slide - preserve template design"""
        content_placeholder = self._find_content_placeholder(slide)
        if content_placeholder:
            try:
                # Format as prominent message
                standout_content = content if content else ["Key Message or Insight"]
                content_list = standout_content if isinstance(standout_content, list) else [standout_content]
                
                # Simple text assignment - let template handle the formatting
                standout_text = '\n'.join(content_list[:2])  # Max 2 items for standout
                content_placeholder.text = standout_text
                self._apply_smart_text_formatting(content_placeholder, standout_text)
                print(f"Set standout slide content")
                
            except Exception as e:
                print(f"Error setting standout content: {e}")
        else:
            print("No content placeholder found for standout slide")

    def _format_any_slide_content(self, slide, content, slide_type):
        """SIMPLIFIED: Format any slide content with table detection and smart formatting"""
        content_placeholder = self._find_content_placeholder(slide)
        
        # SPECIAL HANDLING: For Thank You slide, remove content placeholder to avoid duplication
        if slide_type == "THANK_YOU_SLIDE":
            if content_placeholder:
                try:
                    # Remove the content placeholder entirely - only keep the title
                    slide.shapes._spTree.remove(content_placeholder._element)
                    print(f"Removed content placeholder from Thank You slide - title only")
                except Exception as e:
                    print(f"Error removing Thank You slide content placeholder: {e}")
            return
        
        if content_placeholder and content:
            try:
                content_list = content if isinstance(content, list) else [content]
                
                # Skip table detection for Title and Agenda slides
                should_check_table = slide_type not in ["TITLE_SLIDE", "AGENDA_SLIDE"]
                
                # Check if content should be a table (only for appropriate slide types)
                table_info = self._detect_table_content(content_list) if should_check_table else {"is_table": False}
                
                if table_info["is_table"]:
                    # Remove the content placeholder and create a table instead
                    try:
                        slide.shapes._spTree.remove(content_placeholder._element)
                        print(f"Removed text placeholder to create table")
                    except:
                        print("Could not remove placeholder, creating table anyway")
                    
                    # Create table
                    if self._create_table_slide(slide, table_info):
                        print(f"Successfully created table for {slide_type}")
                        return
                    else:
                        print(f"Table creation failed, falling back to text for {slide_type}")
                
                # Standard text formatting (not a table)
                if len(content_list) == 1:
                    # Single item - could be paragraph or single point
                    content_placeholder.text = str(content_list[0])
                else:
                    # Multiple items - DON'T add manual bullets, let PowerPoint's template handle it
                    content_text = '\n'.join([str(item) for item in content_list[:6]])
                    content_placeholder.text = content_text
                
                # Apply smart text formatting with auto-fit
                self._apply_smart_text_formatting(content_placeholder, content_text if len(content_list) > 1 else str(content_list[0]))
                    
                print(f"Successfully set content for {slide_type}: {len(content_list)} items")
                
            except Exception as e:
                print(f"Error setting {slide_type} content: {e}")
                # Ultimate fallback
                try:
                    simple_text = str(content[0]) if isinstance(content, list) and content else str(content)
                    content_placeholder.text = simple_text
                    self._apply_smart_text_formatting(content_placeholder, simple_text)
                    print("Used simple text fallback")
                except Exception as e2:
                    print(f"All methods failed for {slide_type}: {e2}")
        else:
            print(f"No content placeholder or content for {slide_type}. Placeholder: {bool(content_placeholder)}, Content: {bool(content)}")

    def _apply_title_format_by_type(self, title_shape, slide_type: str):
        """Apply title formatting based on slide type"""
        if title_shape.has_text_frame:
            for paragraph in title_shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    font = run.font
                    font.name = 'Calibri'
                    
                    # Different sizes for different slide types
                    if slide_type in ["TITLE_SLIDE", "THANK_YOU_SLIDE"]:
                        font.size = Pt(36)
                    else:
                        font.size = Pt(32)
                    
                    font.bold = True
                    font.color.rgb = self.PRIMARY_COLOR


    def _find_content_placeholder(self, slide):
        """Find the content placeholder in a slide (NOT the title placeholder)"""
        print(f"Looking for content placeholder in slide with {len(slide.placeholders)} placeholders")
        
        # First, let's inspect what placeholders we actually have
        for i, placeholder in enumerate(slide.placeholders):
            try:
                has_text_frame = hasattr(placeholder, 'has_text_frame') and placeholder.has_text_frame
                is_title = placeholder == slide.shapes.title
                print(f"Placeholder {i}: has_text_frame={has_text_frame}, is_title={is_title}")
            except Exception as e:
                print(f"Placeholder {i}: Error checking - {e}")
        
        # Try to find a content placeholder (NOT the title)
        for i, placeholder in enumerate(slide.placeholders):
            try:
                # CRITICALLY IMPORTANT: Skip the title placeholder completely
                if slide.shapes.title and placeholder == slide.shapes.title:
                    print(f"Skipping placeholder {i} (it's the title placeholder)")
                    continue
                    
                # Check if this placeholder can hold text
                if hasattr(placeholder, 'has_text_frame') and placeholder.has_text_frame:
                    # Additional check - try to access the text_frame
                    _ = placeholder.text_frame
                    print(f"Found working content placeholder at index {i}")
                    return placeholder
                    
            except Exception as e:
                print(f"Placeholder {i}: Error accessing - {e}")
                continue
        
        print("No working content placeholder found!")
        return None

    def _apply_smart_text_formatting(self, content_placeholder, content_text: str):
        """
        Apply smart text formatting with auto-fit and word wrap
        
        HOW TO ADJUST TEXT POSITIONING IN THE FUTURE:
        
        1. MARGINS (text_frame.margin_*):
           - margin_left/right: Controls horizontal spacing from slide edges
           - margin_top/bottom: Controls vertical spacing from slide edges
           - Use Inches(value) - typical range: 0.1 to 0.8 inches
        
        2. VERTICAL ALIGNMENT (text_frame.vertical_anchor):
           - MSO_ANCHOR.TOP: Text starts at top of text box
           - MSO_ANCHOR.MIDDLE: Text centered vertically in text box  
           - MSO_ANCHOR.BOTTOM: Text aligned to bottom of text box
        
        3. HORIZONTAL ALIGNMENT (paragraph.alignment):
           - PP_ALIGN.LEFT: Left-aligned (best for bullets)
           - PP_ALIGN.CENTER: Center-aligned
           - PP_ALIGN.RIGHT: Right-aligned
           - PP_ALIGN.JUSTIFY: Justified text
        
        4. FONT SIZE: Adjust base_font_size values below for different text lengths
        
        5. AUTO-FIT BEHAVIOR (text_frame.auto_size):
           - MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE: Shrinks text to fit
           - MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT: Expands text box
           - MSO_AUTO_SIZE.NONE: No auto-sizing
        """
        if not content_placeholder.has_text_frame:
            return
            
        try:
            text_frame = content_placeholder.text_frame
            
            # Enable word wrap for better text flow
            text_frame.word_wrap = True
            
            # Configure auto-sizing behavior
            from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR
            text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            
            # Set vertical text alignment - center vertically within the text box
            text_frame.vertical_anchor = MSO_ANCHOR.TOP  # Options: TOP, MIDDLE, BOTTOM
            
            # Optimized margins for 16:9 aspect ratio (16" x 9" slide)
            text_frame.margin_left = Inches(0.4)    # Professional left margin
            text_frame.margin_right = Inches(0.4)   # Balanced right margin  
            text_frame.margin_top = Inches(0.3)     # Clean top spacing
            text_frame.margin_bottom = Inches(0.4)  # Adequate bottom spacing
            
            # Set a good starting font size - let auto-fit adjust as needed
            char_count = len(content_text.strip())
            if char_count <= 200:
                base_font_size = 24  # Start larger for short content
            elif char_count <= 400:
                base_font_size = 22  # Medium starting size
            else:
                base_font_size = 20  # Smaller starting size for long content
            
            # Apply base font size and alignment to all content
            from pptx.enum.text import PP_ALIGN
            for paragraph in text_frame.paragraphs:
                # Set paragraph alignment (left is usually best for bullet points)
                paragraph.alignment = PP_ALIGN.LEFT
                
                for run in paragraph.runs:
                    run.font.size = Pt(base_font_size)
                    run.font.name = 'Calibri'
                        
            print(f"Applied auto-fit formatting with {base_font_size}pt base size for {char_count} characters")
            
        except Exception as e:
            print(f"Error applying smart text formatting: {e}")

    def _detect_table_content(self, content_list: list) -> dict:
        """Detect if content should be presented as a table - only when there's comparative/structured data"""
        if not content_list or len(content_list) < 3:  # Require at least 3 items for comparison
            return {"is_table": False}
        
        # Look for patterns that suggest tabular/comparative data
        table_patterns = [
            r'\w+\s*:\s*\$[\d,]+',      # Item: $amount
            r'\w+\s*:\s*\d+\.?\d*%',    # Item: percentage  
            r'\w+\s*:\s*\d+',           # Item: number
            r'[A-Z][^:]*:\s*.+',        # Category: description
        ]
        
        # Also look for comparison indicators
        comparison_keywords = [
            'vs', 'versus', 'compared to', 'difference', 'increase', 'decrease',
            'before', 'after', 'baseline', 'target', 'actual', 'budget',
            'quarter', 'year', 'month', 'period', 'phase'
        ]
        
        matches = 0
        has_comparison_context = False
        
        # Check for structured data patterns
        for item in content_list[:6]:  # Check first 6 items
            item_str = str(item).lower()
            for pattern in table_patterns:
                if re.search(pattern, str(item)):
                    matches += 1
                    break
            
            # Check for comparison context
            if any(keyword in item_str for keyword in comparison_keywords):
                has_comparison_context = True
        
        # Only create table if:
        # 1. 60%+ of items match table patterns (raised from 50%)
        # 2. There are at least 3 matches (comparative data)
        # 3. Content suggests comparison/structured analysis
        if (matches >= len(content_list) * 0.6 and 
            matches >= 3 and 
            (has_comparison_context or matches >= 4)):
            table_data = self._parse_table_data(content_list)
            return {
                "is_table": True,
                "rows": len(table_data),
                "cols": 2,  # Most common: label/value pairs
                "data": table_data
            }
        
        return {"is_table": False}

    def _parse_table_data(self, content_list: list) -> list:
        """Parse content into table rows with proper header detection"""
        table_data = []
        
        # Check if first item looks like headers (no colon, contains common header words)
        first_item = str(content_list[0]).strip() if content_list else ""
        header_indicators = ['phase', 'description', 'item', 'value', 'category', 'type', 'name', 'amount', 'date', 'status']
        
        has_headers = (
            ':' not in first_item and  # Headers typically don't have colons
            any(indicator in first_item.lower() for indicator in header_indicators) and
            len(content_list) > 1
        )
        
        if has_headers:
            # Create proper headers based on content analysis
            sample_content = content_list[1] if len(content_list) > 1 else content_list[0]
            
            if 'phase' in first_item.lower():
                headers = ["Phase", "Description"]
            elif 'budget' in first_item.lower() or '$' in str(sample_content):
                headers = ["Item", "Amount"]
            elif 'team' in first_item.lower() or 'member' in str(sample_content).lower():
                headers = ["Team", "Size"]
            elif 'metric' in first_item.lower() or '%' in str(sample_content):
                headers = ["Metric", "Value"]
            elif 'timeline' in first_item.lower() or any(month in str(sample_content).lower() for month in ['january', 'february', 'march', 'april', 'may', 'june', 'july', 'august', 'september', 'october', 'november', 'december']):
                headers = ["Deliverable", "Timeline"]
            else:
                # Generic headers based on content structure
                headers = ["Category", "Description"]
            
            table_data.append(headers)
            
            # Process all items as data rows (skip first item if it was just header indicator)
            start_index = 1 if any(indicator in first_item.lower() for indicator in header_indicators) else 0
            
            for item in content_list[start_index:]:
                item_str = str(item).strip()
                
                # Split on colon for key:value pairs
                if ':' in item_str:
                    parts = item_str.split(':', 1)
                    if len(parts) == 2:
                        table_data.append([parts[0].strip(), parts[1].strip()])
                    else:
                        table_data.append([item_str, ""])
                else:
                    # If no colon, put entire text in first column
                    table_data.append([item_str, ""])
        else:
            # No headers detected - add generic headers and process as key:value pairs
            # Analyze content to determine appropriate headers
            sample_content = str(content_list[0]) if content_list else ""
            
            if '$' in sample_content:
                headers = ["Item", "Amount"]
            elif '%' in sample_content:
                headers = ["Factor", "Percentage"]
            elif any(month in sample_content.lower() for month in ['january', 'february', 'march', 'april', 'may', 'june', 'july', 'august', 'september', 'october', 'november', 'december', 'q1', 'q2', 'q3', 'q4']):
                headers = ["Activity", "Timeline"]
            elif 'team' in sample_content.lower() or 'department' in sample_content.lower():
                headers = ["Department", "Details"]
            else:
                headers = ["Item", "Details"]
            
            table_data.append(headers)
            
            for item in content_list:
                item_str = str(item).strip()
                
                # Split on colon for key:value pairs
                if ':' in item_str:
                    parts = item_str.split(':', 1)
                    if len(parts) == 2:
                        table_data.append([parts[0].strip(), parts[1].strip()])
                    else:
                        table_data.append([item_str, ""])
                else:
                    # If no colon, put entire text in first column
                    table_data.append([item_str, ""])
        
        return table_data[:8]  # Limit to 8 rows for readability

    def _create_table_slide(self, slide, table_info: dict):
        """Create a table on the slide"""
        try:
            rows = table_info["rows"]
            cols = table_info["cols"] 
            data = table_info["data"]
            
            # Position table optimally for 16:9 aspect ratio (16" x 9" slide)
            left = Inches(0.8)     # Start closer to left edge
            top = Inches(2.2)      # Below title area
            width = Inches(14.4)   # Almost full width (16" - 0.8" left - 0.8" right)
            height = Inches(5.5)   # Good height for 16:9 ratio
            
            # Add table shape
            table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
            table = table_shape.table
            
            # Set column widths optimized for 16:9 ratio and readability
            table.columns[0].width = Inches(6.0)   # First column (labels) - slightly larger
            table.columns[1].width = Inches(8.4)   # Second column (values) - use remaining space
            
            # Fill table with data
            for row_idx, row_data in enumerate(data):
                if row_idx >= rows:
                    break
                    
                for col_idx, cell_value in enumerate(row_data):
                    if col_idx >= cols:
                        break
                        
                    cell = table.cell(row_idx, col_idx)
                    cell.text = str(cell_value)
                    
                    # Format cell text
                    paragraph = cell.text_frame.paragraphs[0]
                    paragraph.alignment = PP_ALIGN.LEFT
                    
                    for run in paragraph.runs:
                        run.font.name = 'Calibri'
                        
                        # Header row formatting (first row)
                        if row_idx == 0:
                            run.font.size = Pt(16)
                            run.font.bold = True
                            run.font.color.rgb = RGBColor(255, 255, 255)  # White text
                            # Set header row background color
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = RGBColor(88, 77, 193)  # Purple theme color
                        else:
                            # Content rows
                            run.font.size = Pt(14)
                            # Make first column bold (labels) for content rows
                            if col_idx == 0:
                                run.font.bold = True
            
            print(f"Created table with {rows} rows and {cols} columns")
            return True
            
        except Exception as e:
            print(f"Error creating table: {e}")
            return False

    def _apply_content_format(self, paragraph):
        """Apply minimal content formatting - preserve template fonts"""
        # Only set the bullet level, don't override template fonts/colors
        paragraph.level = 0
        # Remove all font overrides to preserve template formatting
        # Template should handle font, size, and color