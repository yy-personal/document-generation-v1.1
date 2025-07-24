#!/usr/bin/env python3
"""
Template Analysis Tool - Examines PowerPoint template structure
"""
from pptx import Presentation
import os

def analyze_template(template_path):
    """Analyze PowerPoint template structure for compatibility"""
    print(f"Analyzing template: {template_path}")
    print("=" * 60)
    
    if not os.path.exists(template_path):
        print(f"[ERROR] Template file not found: {template_path}")
        return False
    
    try:
        # Load the template
        prs = Presentation(template_path)
        
        print(f"[OK] Template loaded successfully")
        print(f"[INFO] Slide dimensions: {prs.slide_width.inches:.1f}\" x {prs.slide_height.inches:.1f}\"")
        print(f"[INFO] Number of slide layouts: {len(prs.slide_layouts)}")
        print(f"[INFO] Number of slides in template: {len(prs.slides)}")
        
        # Check if this is 16:9 aspect ratio (more precise check)
        aspect_ratio = prs.slide_width.inches / prs.slide_height.inches
        is_16_9 = abs(aspect_ratio - (16/9)) < 0.1  # Allow small tolerance
        print(f"[INFO] Aspect ratio: {aspect_ratio:.2f} {'(16:9)' if is_16_9 else '(4:3 or other)'}")
        
        # Analyze slide layouts with better detection
        print("\n[LAYOUTS] SLIDE MASTER LAYOUTS AVAILABLE:")
        for i, layout in enumerate(prs.slide_layouts):
            try:
                layout_name = layout.name if hasattr(layout, 'name') else f"Layout {i}"
                print(f"  {i}: {layout_name}")
                
                # Check placeholders in this layout with more detail
                placeholders_info = []
                title_found = False
                content_found = False
                
                for ph_idx, placeholder in enumerate(layout.placeholders):
                    try:
                        ph_type = placeholder.placeholder_format.type
                        
                        # Check for common placeholder types
                        if hasattr(placeholder, 'has_text_frame') and placeholder.has_text_frame:
                            # Try to determine if it's title or content
                            if 'title' in str(ph_type).lower() or ph_idx == 0:
                                placeholders_info.append(f"Title({ph_type})")
                                title_found = True
                            else:
                                placeholders_info.append(f"Content({ph_type})")
                                content_found = True
                        else:
                            placeholders_info.append(f"Other({ph_type})")
                            
                    except Exception as ph_error:
                        placeholders_info.append(f"Unknown({ph_idx})")
                
                if placeholders_info:
                    print(f"     Placeholders: {', '.join(placeholders_info[:6])}")
                else:
                    print(f"     Placeholders: None found")
                
                # Indicate compatibility
                if title_found and content_found:
                    print(f"     [COMPATIBLE] Has title and content areas")
                elif title_found or content_found:
                    print(f"     [PARTIAL] Has {'title' if title_found else 'content'} area only")
                else:
                    print(f"     [LIMITED] No text placeholders detected")
                    
            except Exception as e:
                print(f"  {i}: Layout analysis failed - {e}")
        
        # Check if template has existing slides
        if len(prs.slides) > 0:
            print(f"\n[SLIDES] EXISTING SLIDES IN TEMPLATE:")
            for i, slide in enumerate(prs.slides):
                try:
                    title_text = "No title"
                    if slide.shapes.title:
                        title_text = slide.shapes.title.text or "Empty title"
                    print(f"  Slide {i+1}: {title_text}")
                except:
                    print(f"  Slide {i+1}: Analysis failed")
        
        # Test creating slides from different layouts
        print(f"\n[TEST] COMPATIBILITY TEST:")
        compatible_layouts = []
        
        for i, layout in enumerate(prs.slide_layouts):
            try:
                test_slide = prs.slides.add_slide(layout)
                layout_name = layout.name if hasattr(layout, 'name') else f"Layout {i}"
                
                # Test title capability
                title_works = False
                if test_slide.shapes.title:
                    try:
                        test_slide.shapes.title.text = f"Test Title {i}"
                        title_works = True
                    except:
                        pass
                
                # Test content capability  
                content_placeholders = []
                for placeholder in test_slide.placeholders:
                    try:
                        if (placeholder != test_slide.shapes.title and 
                            hasattr(placeholder, 'has_text_frame') and 
                            placeholder.has_text_frame):
                            placeholder.text = f"Test content {i}"
                            content_placeholders.append(placeholder)
                    except:
                        pass
                
                # Evaluate layout compatibility
                if title_works and len(content_placeholders) > 0:
                    print(f"[OK] Layout {i} ({layout_name}): Full compatibility")
                    compatible_layouts.append((i, layout_name, "full"))
                elif title_works or len(content_placeholders) > 0:
                    print(f"[PARTIAL] Layout {i} ({layout_name}): Partial compatibility")
                    compatible_layouts.append((i, layout_name, "partial"))
                else:
                    print(f"[WARN] Layout {i} ({layout_name}): Limited compatibility")
                    
            except Exception as e:
                print(f"[ERROR] Layout {i} test failed: {e}")
        
        if not compatible_layouts:
            print("[ERROR] No compatible layouts found")
            return False
        else:
            print(f"[OK] Found {len(compatible_layouts)} usable layouts")
        
        print(f"\n[SUMMARY] COMPATIBILITY SUMMARY:")
        print(f"[OK] Template is compatible with python-pptx")
        print(f"[OK] Can be used as base for presentation generation")
        
        # Enhanced recommendations
        print(f"\n[RECOMMENDATIONS]:")
        
        if not is_16_9:
            print(f"[WARN] Template is not 16:9 aspect ratio ({aspect_ratio:.2f})")
            print(f"       Consider using 16:9 for modern presentations")
        else:
            print(f"[OK] Template uses recommended 16:9 aspect ratio")
        
        fully_compatible = len([l for l in compatible_layouts if l[2] == "full"])
        if fully_compatible == 0:
            print(f"[WARN] No layouts have both title and content placeholders")
            print(f"       Add title and content placeholders to your slide masters")
        elif fully_compatible < 2:
            print(f"[WARN] Only {fully_compatible} layout fully compatible")
            print(f"       Consider adding title+content placeholders to more layouts")
        else:
            print(f"[OK] {fully_compatible} layouts are fully compatible")
        
        if len(prs.slide_layouts) < 3:
            print(f"[WARN] Limited slide layouts ({len(prs.slide_layouts)})")
            print(f"       Consider adding more layout variety for different content types")
        else:
            print(f"[OK] Good variety of slide layouts ({len(prs.slide_layouts)})")
        
        # Specific optimization guidance
        print(f"\n[OPTIMIZATION TIPS]:")
        print(f"1. For TITLE slides: Ensure Layout 0 has title placeholder")
        print(f"2. For CONTENT slides: Add both title and content text placeholders")
        print(f"3. For TABLE slides: Use content placeholder - tables will replace text")
        print(f"4. Consider naming layouts descriptively (Title, Content, Two Column, etc.)")
        
        return True
        
    except Exception as e:
        print(f"[ERROR] Failed to analyze template: {e}")
        return False

def main():
    """Analyze the default template"""
    template_path = "templates/default_template.pptx"
    
    print("POWERPOINT TEMPLATE ANALYSIS")
    print("=" * 60)
    
    if analyze_template(template_path):
        print(f"\n[NEXT STEPS]:")
        print(f"1. Enable templates in config.py: set 'use_templates': True")
        print(f"2. Test with: python test_poc.py --table-only")
        print(f"3. Check generated presentations use your template design")
    else:
        print(f"\n[ERROR] Template needs fixes before it can be used")

if __name__ == "__main__":
    main()